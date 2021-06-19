#!/usr/bin/env python
# -*- coding: utf-8 -*-
# @Time    : 2021/5/2 上午9:30
# @Author  : silianpan
# @Site    :
# @File    : file.py
# @Software: PyCharm

from pptx import Presentation
from docx import Document
word_file = Document()
file_path = '/Users/liupan/work/code/seal-python-tools/ppt_process/员工管理制度.pptx'
ppt = Presentation(file_path)
for i in ppt.slides:
    for j in i.shapes:
        if j.has_text_frame:
            text_frame = j.text_frame
            for paragraph in text_frame.paragraphs:
                word_file.add_paragraph(paragraph.text)
save_path = '/Users/liupan/work/code/seal-python-tools/ppt_process/员工管理制度.docx'
word_file.save(save_path)
