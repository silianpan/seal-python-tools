#!/usr/bin/env python
# -*- coding: utf-8 -*-
# @Time    : 2021/5/2 上午9:30
# @Author  : silianpan
# @Site    :
# @File    : file.py
# @Software: PyCharm

from pathlib import Path
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
src_folder = Path('/Users/liupan/work/code/seal-python-tools/ppt_process/公司产品介绍/')
ppt_file = Presentation()
sl_w = Cm(40)
sl_h = Cm(22.5)
ppt_file.slide_width = sl_w
ppt_file.slide_height = sl_h
wb = load_workbook(src_folder / '产品信息表.xlsx')
ws = wb.active
for row in range(2, ws.max_row + 1):
    slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(230, 230, 230)
    rec = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left=0, top=(sl_h - Cm(10)) // 2, width=sl_w, height=Cm(10))
    rec.fill.solid()
    rec.fill.fore_color.rgb = RGBColor(197, 90, 17)
    rec.line.fill.background()
    pic_path = src_folder / (str(ws.cell(row=row, column=1).value) + '.png')
    pic = slide.shapes.add_picture(str(pic_path), left=Cm(2), top=0, width=Cm(15))
    pic.top = (sl_h - pic.height) // 2
    pic.click_action.hyperlink.address = str(ws.cell(row=row, column=5).value)
    txt = slide.shapes.add_textbox(left=sl_w / 2, top=0, width=sl_w / 2, height=sl_h)
    txt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    txt.text_frame.word_wrap = True
    for col in range(2, 5):
        p = txt.text_frame.add_paragraph()
        p.text = f'【{ws.cell(row=1, column=col).value}】{ws.cell(row=row, column=col).value}\n'
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.font.name = '华文中宋'
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
ppt_file.save('/Users/liupan/work/code/seal-python-tools/ppt_process/公司产品介绍.pptx')
